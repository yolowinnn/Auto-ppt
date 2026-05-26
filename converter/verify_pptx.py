#!/usr/bin/env python3
"""Introspect PPTX produced by html2pptx — deeper version showing colors/fonts."""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.dml.color import RGBColor


SHAPE_TYPE_NAMES = {
    1: "AUTO_SHAPE", 3: "CHART", 13: "PICTURE", 14: "PLACEHOLDER",
    17: "TEXT_BOX", 19: "TABLE", 20: "CANVAS", 23: "FREEFORM",
    24: "FREEFORM_BUILD", 28: "LINE", 6: "GROUP",
}


def emu_to_in(emu):
    return Emu(emu).inches


def shape_fill_hex(shape):
    try:
        f = shape.fill
        if f.type is None:
            return "—"
        if f.type == 1:  # solid
            c = f.fore_color.rgb
            return f"#{c}"
        return f"type={f.type}"
    except Exception:
        return "?"


def first_run_props(shape):
    if not shape.has_text_frame:
        return None
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            color_hex = "—"
            try:
                if run.font.color and run.font.color.type is not None:
                    color_hex = f"#{run.font.color.rgb}"
            except Exception:
                pass
            return {
                "font": run.font.name or "—",
                "size": run.font.size.pt if run.font.size else None,
                "bold": run.font.bold,
                "color": color_hex,
            }
    return None


def describe_slide(idx, slide):
    print(f"\n--- Slide {idx + 1} ---")
    counts = {}
    for shape in slide.shapes:
        st = SHAPE_TYPE_NAMES.get(shape.shape_type, f"TYPE_{shape.shape_type}")
        counts[st] = counts.get(st, 0) + 1
    print(f"  Shape inventory: {counts}")

    pic_count = 0
    txt_count = 0
    for i, shape in enumerate(slide.shapes):
        st = SHAPE_TYPE_NAMES.get(shape.shape_type, f"TYPE_{shape.shape_type}")
        x, y, w, h = (emu_to_in(v) for v in (shape.left, shape.top, shape.width, shape.height))
        pos = f"x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f}"
        fill = shape_fill_hex(shape)

        if st == "PICTURE":
            pic_count += 1
            print(f"  [{i:02d}] PICTURE  {pos}  fill={fill}")
        elif not shape.has_text_frame or not shape.text_frame.text.strip():
            print(f"  [{i:02d}] {st}      {pos}  fill={fill}")
        elif shape.has_text_frame and shape.text_frame.text.strip():
            txt_count += 1
            props = first_run_props(shape) or {}
            text = shape.text_frame.text.strip().replace("\n", " | ")[:80]
            font = props.get("font") or "—"
            size = props.get("size")
            bold = "B" if props.get("bold") else " "
            color = props.get("color") or "—"
            print(f"  [{i:02d}] {st}  {pos}  fill={fill}  font={font} size={size} {bold} color={color}")
            print(f"        text=\"{text}\"")

    print(f"  → editable text shapes: {txt_count}, picture (screenshot) shapes: {pic_count}")


def main(pptx_path):
    p = Path(pptx_path)
    if not p.exists():
        print(f"FAIL: {pptx_path} does not exist"); sys.exit(1)
    pres = Presentation(str(p))
    sw = emu_to_in(pres.slide_width)
    sh = emu_to_in(pres.slide_height)
    print(f"PPTX: {pptx_path}")
    print(f"File size: {p.stat().st_size} bytes")
    print(f"Slide canvas: {sw:.2f} x {sh:.2f} in")
    print(f"Total slides: {len(pres.slides)}")
    for i, slide in enumerate(pres.slides):
        describe_slide(i, slide)


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "out/demo.pptx")
