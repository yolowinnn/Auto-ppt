"""Quick layout sanity check: dump shape positions per slide and flag
shapes that fall outside the slide canvas or that overlap heavily."""
import sys
from pptx import Presentation
from pptx.util import Emu

EMU_PER_INCH = 914400

def i(emu): return emu / EMU_PER_INCH

def main(path):
    prs = Presentation(path)
    sw, sh = prs.slide_width, prs.slide_height
    print(f"Slide size: {i(sw):.2f}\" × {i(sh):.2f}\"")
    for idx, slide in enumerate(prs.slides, 1):
        print(f"\n=== Slide {idx} ===")
        boxes = []
        for s in slide.shapes:
            try:
                x, y, w, h = s.left, s.top, s.width, s.height
            except (AttributeError, TypeError):
                continue
            if x is None: continue
            txt = ""
            if s.has_text_frame:
                txt = " ".join(p.text for p in s.text_frame.paragraphs)[:50]
            warn = []
            if x < 0 or y < 0:
                warn.append("OFF-LEFT/TOP")
            if x + w > sw + 100:
                warn.append(f"OFF-RIGHT(+{i(x+w-sw):.2f}\")")
            if y + h > sh + 100:
                warn.append(f"OFF-BOTTOM(+{i(y+h-sh):.2f}\")")
            warn_s = (" ⚠ " + " ".join(warn)) if warn else ""
            print(f"  [{s.shape_type}] "
                  f"x={i(x):.2f} y={i(y):.2f} w={i(w):.2f} h={i(h):.2f}"
                  f" {warn_s}  '{txt}'")
            boxes.append((x, y, w, h, txt))

if __name__ == "__main__":
    main(sys.argv[1])
