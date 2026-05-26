"""
PptxBuilder — reusable layout primitives.

Each public method `add_<layout>_slide(spec)` adds one slide based on a
spec dict. Style-specific behavior (cover layout, footer text, header
chrome) is provided by a `Style` object passed at construction time.

To add a new layout:
  1. Add `add_<name>_slide(self, spec)` here.
  2. Wire it into the dispatch in render.py.
  3. Document required spec fields in SKILL.md and add an example.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

from theme import (
    SLIDE_W, SLIDE_H, MARGIN_X, MARGIN_TOP, TITLE_BAR_Y, CONTENT_TOP, FOOTER_Y,
    GREEN_PRIMARY, GREEN_DARK, GREEN_BRIGHT, GRAY_CALLOUT, GRAY_DIVIDER,
    TEXT_DARK, TEXT_BODY, TEXT_MUTED, WHITE,
    KPI_GOOD, KPI_WARN, KPI_BAD, KPI_INFO, NEG_GRAY, POS_GREEN,
    FONT_LATIN, FONT_CJK, FONT_TITLE_LATIN, FONT_TITLE_CJK,
    SZ_COVER_TITLE, SZ_COVER_SUB, SZ_SECTION, SZ_TITLE, SZ_SUBTITLE,
    SZ_BODY, SZ_SMALL, SZ_FOOTER,
    LOGO_INDUSTRIALMIND, LOGO_TAOMO, has_logo,
)


# ---------- Low-level helpers ------------------------------------------------

def _set_run_font(run, *, size=None, bold=False, color=None,
                  latin=FONT_LATIN, cjk=FONT_CJK):
    """Apply font name (with CJK fallback), size, bold, color to a run."""
    if size is not None:
        run.font.size = size
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    # Remove existing east-asian font if present
    for ea in rPr.findall(qn("a:ea")):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", cjk)


def _add_text(slide, x, y, w, h, text, *, size=SZ_BODY, bold=False,
              color=TEXT_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              latin=FONT_LATIN, cjk=FONT_CJK):
    """Add a textbox; returns the textbox shape."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        if i == 0:
            run = p.add_run()
        else:
            new_p = tf.add_paragraph()
            new_p.alignment = align
            run = new_p.add_run()
        run.text = line
        _set_run_font(run, size=size, bold=bold, color=color,
                      latin=latin, cjk=cjk)
    return tb


def _add_bullet_block(slide, x, y, w, h, items, *,
                      size=SZ_BODY, color=TEXT_BODY,
                      bullet_color=None, line_spacing=1.25,
                      space_after_pt=4):
    """Add a vertically-flowing bullet list as a single textbox.

    `items` is a list; each item can be:
      - a str → single bullet
      - {"text": str, "sub": [str, ...]} → bullet with sub-bullets
      - {"text": str, "bold": True} → bold bullet (e.g. headline-ish bullets)

    Bullet character is "•" with a 2-space indent so wrapped lines align
    visually under the first character of the text (good-enough hanging
    indent without going through pptx XML).
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0);  tf.margin_bottom = Emu(0)

    first = True
    bc = bullet_color or color
    for item in items:
        if isinstance(item, str):
            text, subs, bold = item, [], False
        else:
            text = item.get("text", "")
            subs = item.get("sub", [])
            bold = bool(item.get("bold", False))

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after_pt)
        # bullet glyph
        r1 = p.add_run()
        r1.text = "•  "
        _set_run_font(r1, size=size, bold=True, color=bc)
        r2 = p.add_run()
        r2.text = text
        _set_run_font(r2, size=size, bold=bold, color=color)

        for sub in subs:
            sp = tf.add_paragraph()
            sp.alignment = PP_ALIGN.LEFT
            sp.line_spacing = line_spacing
            sp.space_after = Pt(space_after_pt - 2)
            sr1 = sp.add_run(); sr1.text = "      –  "
            _set_run_font(sr1, size=Pt(size.pt - 2), color=TEXT_MUTED)
            sr2 = sp.add_run(); sr2.text = sub
            _set_run_font(sr2, size=Pt(size.pt - 2), color=TEXT_MUTED)
    return tb


def _add_rect(slide, x, y, w, h, fill=None, line=None, line_width=None):
    """Add a filled rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.shadow.inherit = False
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        if line_width is not None:
            shape.line.width = line_width
    return shape


def _add_line(slide, x1, y1, x2, y2, color=GREEN_PRIMARY, width=Pt(1.0)):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = width
    return line


def _add_image_or_placeholder(slide, path, x, y, w, h, *, label=None):
    """Add image if file exists; else a labeled placeholder rectangle."""
    if path and os.path.isfile(path):
        try:
            slide.shapes.add_picture(path, x, y, w, h)
            return
        except Exception:
            pass
    # placeholder
    _add_rect(slide, x, y, w, h, fill=RGBColor(0xEC, 0xEF, 0xF1),
              line=GRAY_DIVIDER, line_width=Pt(0.75))
    _add_text(slide, x, y, w, h,
              label or (f"[image missing]\n{path}" if path else "[image]"),
              size=SZ_SMALL, color=TEXT_MUTED,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _resolve_path(spec_dir, p):
    if not p:
        return None
    if os.path.isabs(p):
        return p
    return os.path.normpath(os.path.join(spec_dir, p))


# ---------- Builder ----------------------------------------------------------

class PptxBuilder:
    def __init__(self, style, spec_dir="."):
        self.style = style
        self.spec_dir = spec_dir
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self._blank_layout = self.prs.slide_layouts[6]  # truly blank

    def save(self, path):
        os.makedirs(os.path.dirname(os.path.abspath(path)) or ".", exist_ok=True)
        self.prs.save(path)

    def _new_slide(self):
        slide = self.prs.slides.add_slide(self._blank_layout)
        # white default background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self.style.bg_color
        return slide

    # ---- Chrome (header / footer / logo) — delegated to style ----
    def _draw_chrome(self, slide, *, title=None, subtitle=None,
                     show_header=True, show_footer=True, show_logo=True):
        if show_header and title is not None:
            self.style.draw_header(slide, title=title, subtitle=subtitle)
        if show_logo:
            self.style.draw_logo(slide)
        if show_footer:
            self.style.draw_footer(slide)

    # ============================================================
    # LAYOUTS
    # ============================================================

    def add_title_slide(self, spec):
        slide = self._new_slide()
        self.style.draw_title_slide(slide, spec, self.spec_dir)
        return slide

    def add_toc_slide(self, spec):
        """spec: {title, items[]}"""
        slide = self._new_slide()
        self._draw_chrome(slide, title=spec.get("title", "CONTENT"))
        items = spec.get("items", [])
        y = CONTENT_TOP
        for i, item in enumerate(items):
            label = f"{i+1}. {item}" if spec.get("numbered", False) else item
            _add_text(slide, MARGIN_X, y, Inches(11), Inches(0.5), label,
                      size=Pt(18), color=TEXT_BODY)
            y += Inches(0.5)
        return slide

    def add_section_slide(self, spec):
        """spec: {title, subtitle?}"""
        slide = self._new_slide()
        self.style.draw_section_slide(slide, spec)
        return slide

    def add_bullets_slide(self, spec):
        """spec: {title, subtitle?, items[] — strings or {text, sub:[]}}"""
        slide = self._new_slide()
        self._draw_chrome(slide, title=spec["title"])
        y = CONTENT_TOP
        if spec.get("subtitle"):
            _add_text(slide, MARGIN_X, y, SLIDE_W - 2 * MARGIN_X,
                      Inches(0.5), spec["subtitle"], size=SZ_SUBTITLE,
                      bold=True, color=TEXT_DARK)
            y += Inches(0.6)
        _add_bullet_block(
            slide, MARGIN_X, y,
            SLIDE_W - 2 * MARGIN_X, FOOTER_Y - y - Inches(0.2),
            spec.get("items", []),
            size=SZ_BODY, color=TEXT_BODY,
            bullet_color=self.style.bullet_color)
        return slide

    def add_two_col_slide(self, spec):
        """spec: {title, left:{type,...}, right:{type,...}}
        col types: text | bullets | image | table
        """
        slide = self._new_slide()
        self._draw_chrome(slide, title=spec["title"])
        col_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.4)) // 2
        col_h = Inches(5.5)
        y = CONTENT_TOP
        self._draw_col(slide, MARGIN_X, y, col_w, col_h, spec.get("left", {}))
        self._draw_col(slide, MARGIN_X + col_w + Inches(0.4), y, col_w, col_h,
                       spec.get("right", {}))
        return slide

    def _draw_col(self, slide, x, y, w, h, col):
        t = col.get("type", "text")
        if col.get("heading"):
            _add_text(slide, x, y, w, Inches(0.45), col["heading"],
                      size=SZ_SUBTITLE, bold=True,
                      color=self.style.heading_color)
            y += Inches(0.55)
            h -= Inches(0.55)
        if t == "text":
            _add_text(slide, x, y, w, h, col.get("text", ""),
                      size=SZ_BODY, color=TEXT_BODY)
        elif t == "bullets":
            _add_bullet_block(
                slide, x, y, w, h, col.get("items", []),
                size=col.get("size") or SZ_BODY, color=TEXT_BODY,
                bullet_color=self.style.bullet_color)
        elif t == "image":
            cap_h = Inches(0.3) if col.get("caption") else Emu(0)
            _add_image_or_placeholder(slide,
                _resolve_path(self.spec_dir, col.get("path")),
                x, y, w, h - cap_h, label=col.get("caption"))
            if col.get("caption"):
                _add_text(slide, x, y + h - cap_h, w, cap_h,
                          col["caption"], size=SZ_SMALL, color=TEXT_MUTED,
                          align=PP_ALIGN.CENTER)
        elif t == "table":
            self._draw_table_at(slide, x, y, w, h, col)

    def add_image_grid_slide(self, spec):
        """spec: {title, images:[{path, caption?}], cols?: int}"""
        slide = self._new_slide()
        self._draw_chrome(slide, title=spec["title"])
        if spec.get("subtitle"):
            _add_text(slide, MARGIN_X, CONTENT_TOP - Inches(0.05),
                      Inches(12), Inches(0.4), spec["subtitle"],
                      size=SZ_SUBTITLE, bold=True, color=TEXT_DARK)
            top = CONTENT_TOP + Inches(0.5)
        else:
            top = CONTENT_TOP
        imgs = spec.get("images", [])
        n = len(imgs)
        if n == 0:
            return slide
        cols = spec.get("cols") or (1 if n == 1 else 2 if n <= 4 else 3)
        rows = (n + cols - 1) // cols
        gap = Inches(0.2)
        avail_w = SLIDE_W - 2 * MARGIN_X
        avail_h = Inches(5.4)
        cell_w = (avail_w - gap * (cols - 1)) // cols
        cell_h = (avail_h - gap * (rows - 1)) // rows
        for i, img in enumerate(imgs):
            r, c = i // cols, i % cols
            ix = MARGIN_X + c * (cell_w + gap)
            iy = top + r * (cell_h + gap)
            cap = img.get("caption", "")
            cap_h = Inches(0.3) if cap else Emu(0)
            _add_image_or_placeholder(slide,
                _resolve_path(self.spec_dir, img.get("path")),
                ix, iy, cell_w, cell_h - cap_h, label=cap or None)
            if cap:
                _add_text(slide, ix, iy + cell_h - cap_h, cell_w, cap_h,
                          cap, size=SZ_SMALL, color=TEXT_MUTED,
                          align=PP_ALIGN.CENTER)
        return slide

    def add_kpi_cards_slide(self, spec):
        """spec: {title, headline?, cards:[{value, label, status?: good|warn|bad|info}]}"""
        slide = self._new_slide()
        self._draw_chrome(slide, title=spec["title"])
        y = CONTENT_TOP
        if spec.get("headline"):
            _add_text(slide, MARGIN_X, y, Inches(12), Inches(0.5),
                      spec["headline"], size=SZ_SUBTITLE, bold=True,
                      color=TEXT_DARK)
            y += Inches(0.6)
        cards = spec.get("cards", [])
        n = len(cards)
        if n == 0:
            return slide
        gap = Inches(0.18)
        avail = SLIDE_W - 2 * MARGIN_X
        cw = (avail - gap * (n - 1)) // n
        ch = Inches(2.6)
        color_map = {"good": KPI_GOOD, "warn": KPI_WARN,
                     "bad": KPI_BAD, "info": KPI_INFO}
        for i, c in enumerate(cards):
            cx = MARGIN_X + i * (cw + gap)
            accent = color_map.get(c.get("status", "good"), KPI_GOOD)
            # white card with thin border + colored top bar
            _add_rect(slide, cx, y, cw, ch, fill=WHITE,
                      line=GRAY_DIVIDER, line_width=Pt(0.75))
            _add_rect(slide, cx + Inches(0.15), y + Inches(0.2),
                      cw - Inches(0.3), Pt(2.5), fill=accent)
            _add_text(slide, cx, y + Inches(0.7), cw, Inches(1.0),
                      str(c.get("value", "")),
                      size=Pt(36), bold=True, color=accent,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text(slide, cx, y + ch - Inches(0.7), cw, Inches(0.5),
                      str(c.get("label", "")),
                      size=SZ_SMALL, color=TEXT_BODY,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # optional caption
        if spec.get("caption"):
            _add_text(slide, MARGIN_X, y + ch + Inches(0.25), Inches(12),
                      Inches(0.4), spec["caption"],
                      size=SZ_SMALL, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        return slide

    def add_data_table_slide(self, spec):
        """spec: {title, headline?, columns:[str], rows:[[..],..],
                  highlight_col?: int — color last value in row by good/warn/bad}"""
        slide = self._new_slide()
        self._draw_chrome(slide, title=spec["title"])
        y = CONTENT_TOP
        if spec.get("headline"):
            _add_text(slide, MARGIN_X, y, Inches(12), Inches(0.5),
                      spec["headline"], size=SZ_SUBTITLE, bold=True,
                      color=TEXT_DARK)
            y += Inches(0.6)
        self._draw_table_at(slide, MARGIN_X, y,
                            SLIDE_W - 2 * MARGIN_X, Inches(5.0), spec)
        return slide

    def _draw_table_at(self, slide, x, y, w, h, spec):
        cols = spec.get("columns", [])
        rows = spec.get("rows", [])
        if not cols or not rows:
            return
        nrows = len(rows) + 1
        ncols = len(cols)
        table_shape = slide.shapes.add_table(nrows, ncols, x, y, w, h)
        table = table_shape.table
        # header
        for j, c in enumerate(cols):
            cell = table.cell(0, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.style.table_header_bg
            tf = cell.text_frame
            tf.text = ""
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(c)
            _set_run_font(run, size=SZ_SMALL, bold=True,
                          color=self.style.table_header_fg)
        # body
        hl = spec.get("highlight_col")
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = (
                    WHITE if i % 2 == 0 else RGBColor(0xF7, 0xF9, 0xFA))
                tf = cell.text_frame
                tf.text = ""
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = str(val)
                color = TEXT_BODY
                if hl is not None and j == hl:
                    color = self._status_color_from_value(val)
                _set_run_font(run, size=SZ_SMALL, color=color)

    def _status_color_from_value(self, v):
        try:
            s = str(v).strip().rstrip("%")
            f = float(s)
            if f >= 99.5: return KPI_GOOD
            if f >= 95: return KPI_WARN
            return KPI_BAD
        except (ValueError, TypeError):
            return TEXT_BODY

    def add_comparison_slide(self, spec):
        """spec: {title, before:{label,items[]}, after:{label,items[]}, footer_quote?}"""
        slide = self._new_slide()
        self._draw_chrome(slide, title=spec["title"])
        col_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.5)) // 2
        y = CONTENT_TOP
        before = spec.get("before", {})
        after = spec.get("after", {})
        # headers
        _add_rect(slide, MARGIN_X, y, col_w, Inches(0.5), fill=NEG_GRAY)
        _add_text(slide, MARGIN_X, y, col_w, Inches(0.5),
                  before.get("label", "Before"),
                  size=SZ_SUBTITLE, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        rx = MARGIN_X + col_w + Inches(0.5)
        _add_rect(slide, rx, y, col_w, Inches(0.5), fill=POS_GREEN)
        _add_text(slide, rx, y, col_w, Inches(0.5),
                  after.get("label", "After"),
                  size=SZ_SUBTITLE, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # rows
        cy = y + Inches(0.7)
        b_items = before.get("items", [])
        a_items = after.get("items", [])
        n = max(len(b_items), len(a_items))
        row_h = Inches(0.55)
        for i in range(n):
            if i < len(b_items):
                _add_text(slide, MARGIN_X, cy, Inches(0.5), row_h,
                          "✗", size=SZ_TITLE, bold=True, color=KPI_BAD,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                _add_text(slide, MARGIN_X + Inches(0.5), cy,
                          col_w - Inches(0.5), row_h, b_items[i],
                          size=SZ_BODY, color=TEXT_BODY,
                          anchor=MSO_ANCHOR.MIDDLE)
            if i < len(a_items):
                _add_text(slide, rx, cy, Inches(0.5), row_h,
                          "✓", size=SZ_TITLE, bold=True, color=POS_GREEN,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                _add_text(slide, rx + Inches(0.5), cy,
                          col_w - Inches(0.5), row_h, a_items[i],
                          size=SZ_BODY, color=TEXT_BODY,
                          anchor=MSO_ANCHOR.MIDDLE)
            cy += row_h
        if spec.get("footer_quote"):
            _add_rect(slide, MARGIN_X, Inches(6.5),
                      SLIDE_W - 2 * MARGIN_X, Inches(0.55),
                      fill=GRAY_CALLOUT)
            _add_text(slide, MARGIN_X, Inches(6.5),
                      SLIDE_W - 2 * MARGIN_X, Inches(0.55),
                      spec["footer_quote"],
                      size=SZ_SMALL, bold=True, color=TEXT_DARK,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return slide

    def add_comparison_cards_slide(self, spec):
        """Manus-style two-column card comparison.

        Each column is a white card with a rounded-corner design:
          • colored header (rounded top, flat bottom via a cover strip)
          • white body with ✗ / ✓ icons and item text
          • bottom: light-gray pill footer_quote

        Implementation detail — matching header to card corners:
          1. White ROUNDED_RECTANGLE (full card, light border)
          2. Colored ROUNDED_RECTANGLE (header, adj tuned for same corner radius)
          3. Plain RECTANGLE cover strip (hides header's bottom curve → flat seam)
          4. Header label text on top

        spec: {title, before:{label,items[]}, after:{label,items[]}, footer_quote?}
        """
        from pptx.dml.color import RGBColor as _RGB
        slide = self._new_slide()
        self._draw_chrome(slide, title=spec["title"])

        HEADER_H = Inches(0.60)
        CARD_R   = Inches(0.15)   # corner radius (EMU); used for geom + positioning
        COL_GAP  = Inches(0.35)
        col_w    = (SLIDE_W - 2 * MARGIN_X - COL_GAP) // 2
        card_y   = CONTENT_TOP + Inches(0.05)
        fq_h     = Inches(0.55) if spec.get("footer_quote") else Emu(0)
        fq_gap   = Inches(0.15) if spec.get("footer_quote") else Emu(0)
        card_h   = FOOTER_Y - card_y - fq_h - fq_gap - Inches(0.15)

        # adjustment fractions for python-pptx (0..1 maps to 0..50000 in ooxml)
        CARD_ADJ = CARD_R / min(col_w, card_h)
        HDR_ADJ  = min(CARD_R / min(col_w, HEADER_H), 0.499)

        before  = spec.get("before", {})
        after   = spec.get("after", {})
        items_b = before.get("items", [])
        items_a = after.get("items", [])
        n       = max(len(items_b), len(items_a))
        item_h  = (card_h - HEADER_H - Inches(0.2)) / max(n, 1)
        ICON_W  = Inches(0.50)
        PAD     = Inches(0.18)

        for side in range(2):
            x          = MARGIN_X + side * (col_w + COL_GAP)
            hdr_color  = NEG_GRAY  if side == 0 else POS_GREEN
            label      = (before if side == 0 else after).get("label", "")
            items      = items_b  if side == 0 else items_a
            icon       = "✗"      if side == 0 else "✓"
            icon_color = KPI_BAD  if side == 0 else POS_GREEN

            # ── 1. white card background (rounded) ──────────────────────────
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, col_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = _RGB(0xD0, 0xD5, 0xD8)
            card.line.width = Pt(0.75)
            try:
                card.adjustments[0] = CARD_ADJ
            except Exception:
                pass

            # ── 2. colored header (rounded rect, corners match card) ────────
            hdr = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, col_w, HEADER_H)
            hdr.fill.solid()
            hdr.fill.fore_color.rgb = hdr_color
            hdr.line.fill.background()
            try:
                hdr.adjustments[0] = HDR_ADJ
            except Exception:
                pass

            # ── 3. cover strip → flat bottom seam for header ────────────────
            cover = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, card_y + HEADER_H - CARD_R, col_w, CARD_R + Pt(1))
            cover.fill.solid()
            cover.fill.fore_color.rgb = hdr_color
            cover.line.fill.background()

            # ── 4. header label text ─────────────────────────────────────────
            _add_text(slide, x, card_y, col_w, HEADER_H, label,
                      size=SZ_SUBTITLE, bold=True, color=WHITE,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            # ── 5. item rows ─────────────────────────────────────────────────
            iy = card_y + HEADER_H + Inches(0.10)
            for item in items:
                _add_text(slide, x + PAD, iy, ICON_W, item_h,
                          icon, size=Pt(16), bold=True, color=icon_color,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                _add_text(slide, x + PAD + ICON_W, iy,
                          col_w - ICON_W - PAD * 2, item_h,
                          item, size=SZ_BODY, color=TEXT_BODY,
                          anchor=MSO_ANCHOR.MIDDLE)
                iy += item_h

        # ── footer pill ──────────────────────────────────────────────────────
        if spec.get("footer_quote"):
            fq_y = FOOTER_Y - fq_h - Inches(0.10)
            fq = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                MARGIN_X, fq_y, SLIDE_W - 2 * MARGIN_X, fq_h)
            fq.fill.solid()
            fq.fill.fore_color.rgb = _RGB(0xEE, 0xF1, 0xF3)
            fq.line.fill.background()
            try:
                fq.adjustments[0] = 0.499   # pill-shaped ends
            except Exception:
                pass
            _add_text(slide, MARGIN_X, fq_y,
                      SLIDE_W - 2 * MARGIN_X, fq_h,
                      spec["footer_quote"],
                      size=SZ_SMALL, color=TEXT_BODY,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        return slide

    def add_flow_slide(self, spec):
        """spec: {title, boxes:[{label, color?: green|orange|gray}]}"""
        slide = self._new_slide()
        self._draw_chrome(slide, title=spec["title"])
        boxes = spec.get("boxes", [])
        if not boxes:
            return slide
        n = len(boxes)
        y = Inches(3.2)
        gap = Inches(0.3)
        avail = SLIDE_W - 2 * MARGIN_X
        bw = (avail - gap * (n - 1)) // n
        bh = Inches(1.0)
        color_map = {"green": GREEN_PRIMARY, "orange": RGBColor(0xEF, 0x6C, 0x00),
                     "gray": NEG_GRAY, "blue": KPI_INFO}
        for i, b in enumerate(boxes):
            bx = MARGIN_X + i * (bw + gap)
            color = color_map.get(b.get("color", "green"), GREEN_PRIMARY)
            _add_rect(slide, bx, y, bw, bh, fill=color)
            _add_text(slide, bx, y, bw, bh, b.get("label", ""),
                      size=SZ_BODY, bold=True, color=WHITE,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            if i < n - 1:
                # arrow
                ax = bx + bw
                ay = y + bh / 2
                arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                             ax, ay - Inches(0.1),
                                             gap, Inches(0.2))
                arr.fill.solid()
                arr.fill.fore_color.rgb = TEXT_MUTED
                arr.line.fill.background()
        return slide

    def add_thanks_slide(self, spec):
        slide = self._new_slide()
        self.style.draw_thanks_slide(slide, spec, self.spec_dir)
        return slide

    # ---- New visual layouts -----------------------------------------------

    def add_chart_bar_slide(self, spec):
        """Native bar chart.
        spec: {title, headline?, categories:[str], series:[{name, values:[..]}],
               horizontal?: bool=False, caption?: str}
        """
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

        slide = self._new_slide()
        self._draw_chrome(slide, title=spec["title"])
        y = CONTENT_TOP
        if spec.get("headline"):
            _add_text(slide, MARGIN_X, y, SLIDE_W - 2 * MARGIN_X,
                      Inches(0.5), spec["headline"],
                      size=SZ_SUBTITLE, bold=True, color=TEXT_DARK)
            y += Inches(0.55)
        cdata = CategoryChartData()
        cdata.categories = spec.get("categories", [])
        for s in spec.get("series", []):
            cdata.add_series(s.get("name", ""), s.get("values", []))
        ctype = (XL_CHART_TYPE.BAR_CLUSTERED if spec.get("horizontal")
                 else XL_CHART_TYPE.COLUMN_CLUSTERED)
        h_avail = FOOTER_Y - y - Inches(0.6) if spec.get("caption") else \
                  FOOTER_Y - y - Inches(0.2)
        chart_shape = slide.shapes.add_chart(
            ctype, MARGIN_X, y,
            SLIDE_W - 2 * MARGIN_X, h_avail, cdata)
        chart = chart_shape.chart
        chart.has_title = False
        if len(spec.get("series", [])) > 1:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        else:
            chart.has_legend = False
        # data labels
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = SZ_SMALL
        if spec.get("caption"):
            _add_text(slide, MARGIN_X, FOOTER_Y - Inches(0.4),
                      SLIDE_W - 2 * MARGIN_X, Inches(0.3),
                      spec["caption"], size=SZ_SMALL, color=TEXT_MUTED,
                      align=PP_ALIGN.CENTER)
        return slide

    def add_image_with_callouts_slide(self, spec):
        """Title + left image + right bullet callouts.
        spec: {title, headline?, image:{path, caption?},
               callouts:[str|{text,sub:[]}], image_ratio?: 0.5..0.7}
        """
        slide = self._new_slide()
        self._draw_chrome(slide, title=spec["title"])
        y = CONTENT_TOP
        if spec.get("headline"):
            _add_text(slide, MARGIN_X, y, SLIDE_W - 2 * MARGIN_X,
                      Inches(0.5), spec["headline"],
                      size=SZ_SUBTITLE, bold=True, color=TEXT_DARK)
            y += Inches(0.55)
        ratio = spec.get("image_ratio", 0.6)
        avail = SLIDE_W - 2 * MARGIN_X - Inches(0.4)
        img_w = int(avail * ratio)
        col_w = avail - img_w
        h = FOOTER_Y - y - Inches(0.2)
        img = spec.get("image", {}) or {}
        cap_h = Inches(0.3) if img.get("caption") else Emu(0)
        # image with subtle border
        _add_rect(slide, MARGIN_X - Pt(0.5), y - Pt(0.5),
                  img_w + Pt(1), h - cap_h + Pt(1),
                  fill=None, line=GRAY_DIVIDER, line_width=Pt(0.5))
        _add_image_or_placeholder(
            slide, _resolve_path(self.spec_dir, img.get("path")),
            MARGIN_X, y, img_w, h - cap_h, label=img.get("caption"))
        if img.get("caption"):
            _add_text(slide, MARGIN_X, y + h - cap_h, img_w, cap_h,
                      img["caption"], size=SZ_SMALL, color=TEXT_MUTED,
                      align=PP_ALIGN.CENTER)
        # right callouts
        cx = MARGIN_X + img_w + Inches(0.4)
        _add_bullet_block(slide, cx, y, col_w, h,
                          spec.get("callouts", []),
                          size=SZ_BODY, color=TEXT_BODY,
                          bullet_color=self.style.bullet_color)
        return slide

    def add_metrics_strip_slide(self, spec):
        """A taller hero metrics strip — fewer, larger numbers, with
        optional supporting body below.
        spec: {title, headline?, metrics:[{value, label, sub?}],
               body?: [str|{text,sub}], caption?}
        """
        slide = self._new_slide()
        self._draw_chrome(slide, title=spec["title"])
        y = CONTENT_TOP
        if spec.get("headline"):
            _add_text(slide, MARGIN_X, y, SLIDE_W - 2 * MARGIN_X,
                      Inches(0.5), spec["headline"],
                      size=SZ_SUBTITLE, bold=True, color=TEXT_DARK)
            y += Inches(0.55)
        metrics = spec.get("metrics", [])
        n = len(metrics)
        if n:
            gap = Inches(0.25)
            avail = SLIDE_W - 2 * MARGIN_X
            cw = (avail - gap * (n - 1)) // n
            ch = Inches(1.7)
            for i, m in enumerate(metrics):
                cx = MARGIN_X + i * (cw + gap)
                _add_rect(slide, cx, y, cw, ch, fill=GRAY_CALLOUT)
                _add_rect(slide, cx, y, Pt(3), ch,
                          fill=self.style.accent_color)
                _add_text(slide, cx + Inches(0.15), y + Inches(0.15),
                          cw - Inches(0.3), Inches(0.7),
                          str(m.get("value", "")),
                          size=Pt(32), bold=True,
                          color=self.style.accent_color)
                _add_text(slide, cx + Inches(0.15), y + Inches(0.9),
                          cw - Inches(0.3), Inches(0.35),
                          str(m.get("label", "")),
                          size=Pt(12), bold=True, color=TEXT_DARK)
                if m.get("sub"):
                    _add_text(slide, cx + Inches(0.15),
                              y + Inches(1.25), cw - Inches(0.3),
                              Inches(0.4), str(m["sub"]),
                              size=Pt(10), color=TEXT_MUTED)
            y += ch + Inches(0.3)
        if spec.get("body"):
            body_bottom_pad = Inches(0.6) if spec.get("caption") else Inches(0.2)
            _add_bullet_block(slide, MARGIN_X, y,
                              SLIDE_W - 2 * MARGIN_X,
                              FOOTER_Y - y - body_bottom_pad,
                              spec["body"],
                              size=SZ_BODY, color=TEXT_BODY,
                              bullet_color=self.style.bullet_color)
        if spec.get("caption"):
            _add_text(slide, MARGIN_X, FOOTER_Y - Inches(0.4),
                      SLIDE_W - 2 * MARGIN_X, Inches(0.3),
                      spec["caption"], size=SZ_SMALL, color=TEXT_MUTED,
                      align=PP_ALIGN.CENTER)
        return slide

    def add_three_line_table_slide(self, spec):
        """三线表 (three-line table) — academic-style, no vertical borders.
        Top/bottom rules are thick (2 pt, GREEN_DARK).
        Header-bottom rule is thin (0.75 pt, TEXT_DARK).
        Odd data rows get a very light tint.

        spec: {
          title,
          headline?: str,
          columns: [str | {text, width}],  # width = proportion 0..1
          rows: [[str, ...], ...],
          caption?: str
        }
        Column `width` values should sum to <= 1.0.  Any columns with no
        explicit width share the remaining proportion equally.
        """
        from pptx.dml.color import RGBColor as _RGB
        slide = self._new_slide()
        self._draw_chrome(slide, title=spec["title"])
        y = CONTENT_TOP
        if spec.get("headline"):
            _add_text(slide, MARGIN_X, y, SLIDE_W - 2 * MARGIN_X,
                      Inches(0.5), spec["headline"],
                      size=SZ_SUBTITLE, bold=True, color=TEXT_DARK)
            y += Inches(0.55)

        cols_spec = spec.get("columns", [])
        rows_data = spec.get("rows", [])
        if not cols_spec or not rows_data:
            return slide

        avail_w = SLIDE_W - 2 * MARGIN_X

        # ── parse column specs ──────────────────────────────────────────────
        col_texts, col_props = [], []
        for c in cols_spec:
            if isinstance(c, str):
                col_texts.append(c)
                col_props.append(None)
            else:
                col_texts.append(c.get("text", ""))
                col_props.append(c.get("width"))   # float 0..1 or None

        known = sum(p for p in col_props if p is not None)
        none_n = sum(1 for p in col_props if p is None)
        fill = (1.0 - known) / max(none_n, 1)
        col_props = [p if p is not None else fill for p in col_props]
        col_widths = [int(avail_w * p) for p in col_props]

        HEADER_H = Inches(0.42)
        ROW_H    = Inches(spec.get("row_height", 0.44))
        PAD_X    = Inches(0.12)

        # ── top thick rule ──────────────────────────────────────────────────
        _add_line(slide, MARGIN_X, y, MARGIN_X + avail_w, y,
                  color=GREEN_DARK, width=Pt(2.0))

        # ── header row ──────────────────────────────────────────────────────
        x = MARGIN_X
        for j, col_text in enumerate(col_texts):
            cw = col_widths[j]
            _add_text(slide, x + PAD_X, y, cw - 2 * PAD_X, HEADER_H,
                      col_text, size=SZ_SMALL, bold=True, color=TEXT_DARK,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            x += cw

        # ── header-bottom thin rule ─────────────────────────────────────────
        _add_line(slide, MARGIN_X, y + HEADER_H,
                  MARGIN_X + avail_w, y + HEADER_H,
                  color=TEXT_DARK, width=Pt(0.75))

        # ── body rows ───────────────────────────────────────────────────────
        for i, row in enumerate(rows_data):
            ry = y + HEADER_H + i * ROW_H
            if i % 2 == 1:
                _add_rect(slide, MARGIN_X, ry, avail_w, ROW_H,
                          fill=_RGB(0xF5, 0xF7, 0xF8))
            x = MARGIN_X
            for j, val in enumerate(row):
                cw = col_widths[j]
                is_first = (j == 0)
                _add_text(
                    slide, x + PAD_X, ry, cw - 2 * PAD_X, ROW_H,
                    str(val),
                    size=SZ_SMALL,
                    bold=is_first,
                    color=self.style.accent_color if is_first else TEXT_BODY,
                    align=PP_ALIGN.CENTER if is_first else PP_ALIGN.LEFT,
                    anchor=MSO_ANCHOR.MIDDLE,
                )
                x += cw

        # ── bottom thick rule ───────────────────────────────────────────────
        table_bottom = y + HEADER_H + len(rows_data) * ROW_H
        _add_line(slide, MARGIN_X, table_bottom,
                  MARGIN_X + avail_w, table_bottom,
                  color=GREEN_DARK, width=Pt(2.0))

        # ── caption ─────────────────────────────────────────────────────────
        if spec.get("caption"):
            _add_text(slide, MARGIN_X, table_bottom + Inches(0.18),
                      avail_w, Inches(0.4), spec["caption"],
                      size=SZ_SMALL, color=TEXT_MUTED, align=PP_ALIGN.LEFT)

        return slide

    def add_video_slide(self, spec):
        """Centered video on chrome slide.
        spec: {title, headline?, video:{path, poster?, caption?, aspect?}}
        - poster: still image shown before play; auto-generated frame works
        - aspect: width/height ratio of video (default 16/9 = 1.778)
        - caption: small text below video
        """
        slide = self._new_slide()
        self._draw_chrome(slide, title=spec["title"])
        y = CONTENT_TOP
        if spec.get("headline"):
            _add_text(slide, MARGIN_X, y, SLIDE_W - 2 * MARGIN_X,
                      Inches(0.5), spec["headline"],
                      size=SZ_SUBTITLE, bold=True, color=TEXT_DARK,
                      align=PP_ALIGN.CENTER)
            y += Inches(0.55)
        video = spec.get("video", {}) or {}
        vid_path = _resolve_path(self.spec_dir, video.get("path"))
        poster_path = _resolve_path(self.spec_dir, video.get("poster"))
        aspect = float(video.get("aspect", 16/9))
        caption = video.get("caption")
        cap_h = Inches(0.35) if caption else Emu(0)
        avail_w = SLIDE_W - 2 * MARGIN_X
        avail_h = FOOTER_Y - y - Inches(0.20) - cap_h
        # fit by height first, then by width
        if avail_w / avail_h > aspect:
            vh = avail_h
            vw = int(vh * aspect)
        else:
            vw = avail_w
            vh = int(vw / aspect)
        vx = (SLIDE_W - vw) // 2
        vy = y + (avail_h - vh) // 2
        # subtle border
        _add_rect(slide, vx - Pt(0.5), vy - Pt(0.5),
                  vw + Pt(1), vh + Pt(1),
                  fill=None, line=GRAY_DIVIDER, line_width=Pt(0.5))
        if vid_path and os.path.exists(vid_path):
            slide.shapes.add_movie(
                vid_path, vx, vy, vw, vh,
                poster_frame_image=(poster_path
                                    if poster_path and os.path.exists(poster_path)
                                    else None),
                mime_type="video/mp4",
            )
        else:
            _add_rect(slide, vx, vy, vw, vh, fill=GRAY_CALLOUT)
            _add_text(slide, vx, vy, vw, vh,
                      f"[VIDEO MISSING: {vid_path}]",
                      size=SZ_BODY, color=TEXT_MUTED,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if caption:
            _add_text(slide, MARGIN_X, FOOTER_Y - cap_h - Inches(0.15),
                      avail_w, cap_h, caption,
                      size=SZ_SMALL, color=TEXT_MUTED,
                      align=PP_ALIGN.CENTER)
        return slide
