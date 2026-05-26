"""
Style: pitch — IndustrialMind external pitch deck.

Visual signature (from Data/Intelligent_Assembly_Action_Anomaly_Detection_System.pdf):
  - Cover & thanks: full dark-green background, white logo + white text
  - Content: white background, dark-green title + thin green underline
  - Right-top: small IndustrialMind.ai green logo + wordmark
  - Footer: "© <year> IndustrialMind AI | All Rights Reserved." centered, gray
"""
import os, sys
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from theme import (
    SLIDE_W, SLIDE_H, MARGIN_X, TITLE_BAR_Y, FOOTER_Y,
    GREEN_PRIMARY, GREEN_DARK, GREEN_BRIGHT, GRAY_DIVIDER,
    TEXT_DARK, TEXT_BODY, TEXT_MUTED, WHITE,
    SZ_TITLE, SZ_SUBTITLE, SZ_BODY, SZ_FOOTER, SZ_COVER_TITLE, SZ_COVER_SUB,
    SZ_SECTION,
    LOGO_INDUSTRIALMIND, has_logo,
)
from builder import _add_text, _add_rect, _add_line, _add_image_or_placeholder, _resolve_path


class PitchStyle:
    name = "pitch"
    bg_color = WHITE
    table_header_bg = GREEN_DARK
    table_header_fg = WHITE
    accent_color = GREEN_DARK
    bullet_color = GREEN_PRIMARY
    heading_color = GREEN_DARK
    default_footer = "© 2026 IndustrialMind AI | All Rights Reserved."

    def __init__(self, meta):
        self.meta = meta or {}
        self.footer_text = self.meta.get("footer", self.default_footer)

    def draw_header(self, slide, *, title, subtitle=None):
        _add_text(slide, MARGIN_X, TITLE_BAR_Y - Inches(0.05),
                  Inches(11), Inches(0.5), title,
                  size=SZ_TITLE, bold=True, color=GREEN_DARK)
        # underline
        _add_line(slide, MARGIN_X, TITLE_BAR_Y + Inches(0.5),
                  SLIDE_W - MARGIN_X, TITLE_BAR_Y + Inches(0.5),
                  color=GREEN_PRIMARY, width=Pt(0.75))
        if subtitle:
            _add_text(slide, MARGIN_X, TITLE_BAR_Y + Inches(0.55),
                      Inches(11), Inches(0.4), subtitle,
                      size=SZ_SUBTITLE, color=TEXT_MUTED)

    def draw_logo(self, slide):
        x = SLIDE_W - Inches(1.6)
        y = Inches(0.3)
        if has_logo(LOGO_INDUSTRIALMIND):
            slide.shapes.add_picture(LOGO_INDUSTRIALMIND, x, y,
                                     width=Inches(1.2))
        else:
            self._draw_text_logo(slide, x, y, GREEN_DARK)

    def _draw_text_logo(self, slide, x, y, color, white=False):
        for i, off in enumerate([0.0, 0.18, 0.36]):
            tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                         x + Inches(off), y,
                                         Inches(0.22), Inches(0.32))
            tri.fill.solid()
            tri.fill.fore_color.rgb = color
            tri.line.fill.background()
        _add_text(slide, x - Inches(0.1), y + Inches(0.34), Inches(0.78),
                  Inches(0.22), "IndustrialMind.ai",
                  size=Pt(8), bold=True, color=color, align=PP_ALIGN.CENTER)

    def draw_footer(self, slide):
        _add_text(slide, Inches(0), FOOTER_Y, SLIDE_W, Inches(0.3),
                  self.footer_text,
                  size=SZ_FOOTER, color=TEXT_MUTED,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ---- Cover (full dark green) ----
    def draw_title_slide(self, slide, spec, spec_dir):
        # Override bg to dark green
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = GREEN_DARK
        # logo (white)
        if has_logo(LOGO_INDUSTRIALMIND):
            slide.shapes.add_picture(LOGO_INDUSTRIALMIND,
                                     SLIDE_W // 2 - Inches(0.6),
                                     Inches(1.6), width=Inches(1.2))
        else:
            self._draw_text_logo(slide, SLIDE_W // 2 - Inches(0.27),
                                 Inches(1.6), WHITE)
        # title
        _add_text(slide, MARGIN_X, Inches(2.7), SLIDE_W - 2 * MARGIN_X,
                  Inches(1.2), spec.get("title", ""),
                  size=SZ_COVER_TITLE, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)
        if spec.get("subtitle"):
            _add_text(slide, MARGIN_X, Inches(3.9), SLIDE_W - 2 * MARGIN_X,
                      Inches(0.6), spec["subtitle"],
                      size=Pt(18), color=WHITE, align=PP_ALIGN.CENTER)
        # divider
        _add_line(slide, Inches(2.5), Inches(4.7),
                  SLIDE_W - Inches(2.5), Inches(4.7),
                  color=WHITE, width=Pt(0.5))
        if spec.get("author"):
            _add_text(slide, MARGIN_X, Inches(4.85), SLIDE_W - 2 * MARGIN_X,
                      Inches(0.5), spec["author"],
                      size=SZ_SUBTITLE, bold=True, color=WHITE,
                      align=PP_ALIGN.CENTER)
        if spec.get("affiliation"):
            _add_text(slide, MARGIN_X, Inches(5.35), SLIDE_W - 2 * MARGIN_X,
                      Inches(0.4), spec["affiliation"],
                      size=Pt(12), color=WHITE, align=PP_ALIGN.CENTER)
        _add_text(slide, Inches(0), Inches(7.0), SLIDE_W, Inches(0.3),
                  self.footer_text, size=SZ_FOOTER, color=WHITE,
                  align=PP_ALIGN.CENTER)

    def draw_section_slide(self, slide, spec):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = GREEN_DARK
        _add_text(slide, Inches(0), Inches(3.0), SLIDE_W, Inches(1.2),
                  spec.get("title", "SECTION"),
                  size=SZ_SECTION, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if spec.get("subtitle"):
            _add_text(slide, Inches(0), Inches(4.3), SLIDE_W, Inches(0.6),
                      spec["subtitle"],
                      size=SZ_SUBTITLE, color=WHITE,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    def draw_thanks_slide(self, slide, spec, spec_dir):
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = GREEN_DARK
        if has_logo(LOGO_INDUSTRIALMIND):
            slide.shapes.add_picture(LOGO_INDUSTRIALMIND,
                                     SLIDE_W // 2 - Inches(0.6),
                                     Inches(1.4), width=Inches(1.2))
        else:
            self._draw_text_logo(slide, SLIDE_W // 2 - Inches(0.27),
                                 Inches(1.4), WHITE)
        _add_text(slide, MARGIN_X, Inches(2.5), SLIDE_W - 2 * MARGIN_X,
                  Inches(1.0), spec.get("title", "Thank you"),
                  size=SZ_COVER_TITLE, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER)
        if spec.get("subtitle"):
            _add_text(slide, MARGIN_X, Inches(3.6), SLIDE_W - 2 * MARGIN_X,
                      Inches(0.5), spec["subtitle"],
                      size=SZ_SUBTITLE, color=WHITE, align=PP_ALIGN.CENTER)
        if spec.get("tagline"):
            _add_text(slide, MARGIN_X, Inches(4.2), SLIDE_W - 2 * MARGIN_X,
                      Inches(0.4), spec["tagline"],
                      size=Pt(13), color=WHITE, align=PP_ALIGN.CENTER)
        _add_line(slide, Inches(2.5), Inches(5.0),
                  SLIDE_W - Inches(2.5), Inches(5.0),
                  color=WHITE, width=Pt(0.5))
        cy = Inches(5.2)
        for line in spec.get("contact", []):
            _add_text(slide, MARGIN_X, cy, SLIDE_W - 2 * MARGIN_X,
                      Inches(0.4), line,
                      size=Pt(12), color=WHITE, align=PP_ALIGN.CENTER)
            cy += Inches(0.4)
        _add_text(slide, Inches(0), Inches(7.0), SLIDE_W, Inches(0.3),
                  self.footer_text, size=SZ_FOOTER, color=WHITE,
                  align=PP_ALIGN.CENTER)
