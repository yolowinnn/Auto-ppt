"""
Style: internal — Siemens / TaomoAI weekly tracking deck.

Visual signature (from Data/Simens PPAI Weekly Tracking - 0112.pdf):
  - White background
  - Top-left bold black title, thin green underline (full width)
  - Top-right small green logo (▲ stacked peaks)
  - Bottom centered footer: "Copyright @ TaomoAI 2024. Business Confidential"
  - Cover: split layout — left half text, right half full-bleed cityscape image
  - Thanks: same split layout with "Thanks" + tagline
"""
import os, sys
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from theme import (
    SLIDE_W, SLIDE_H, MARGIN_X, TITLE_BAR_Y, FOOTER_Y,
    GREEN_PRIMARY, GREEN_DARK, GRAY_DIVIDER, TEXT_DARK, TEXT_BODY, TEXT_MUTED,
    WHITE, FONT_LATIN, FONT_CJK, SZ_TITLE, SZ_SUBTITLE, SZ_BODY, SZ_FOOTER,
    SZ_COVER_TITLE, SZ_COVER_SUB, SZ_SECTION,
    LOGO_INDUSTRIALMIND, ASSETS_DIR, has_logo,
)
from builder import _add_text, _add_rect, _add_line, _add_image_or_placeholder, _resolve_path


class InternalStyle:
    name = "internal"
    bg_color = WHITE
    table_header_bg = RGBColor(0xEC, 0xEF, 0xF1)
    table_header_fg = TEXT_DARK
    accent_color = GREEN_PRIMARY
    bullet_color = GREEN_PRIMARY
    heading_color = TEXT_DARK
    default_footer = "Copyright @ TaomoAI 2026. Business Confidential"

    def __init__(self, meta):
        self.meta = meta or {}
        self.footer_text = self.meta.get("footer", self.default_footer)

    # ---- Chrome ----
    def draw_header(self, slide, *, title, subtitle=None):
        # title text (top row)
        _add_text(slide, MARGIN_X, TITLE_BAR_Y - Inches(0.08),
                  SLIDE_W - 2 * MARGIN_X - Inches(1.5), Inches(0.5), title,
                  size=SZ_TITLE, bold=True, color=TEXT_DARK)
        # full-width green underline BELOW the title — never overlaps
        _add_line(slide, MARGIN_X, TITLE_BAR_Y + Inches(0.55),
                  SLIDE_W - MARGIN_X - Inches(1.5),
                  TITLE_BAR_Y + Inches(0.55),
                  color=GREEN_PRIMARY, width=Pt(1.25))

    def draw_logo(self, slide):
        # right-top corner; ~1in wide
        x = SLIDE_W - Inches(1.4)
        y = Inches(0.25)
        if has_logo(LOGO_INDUSTRIALMIND):
            slide.shapes.add_picture(LOGO_INDUSTRIALMIND, x, y,
                                     width=Inches(1.0))
        else:
            self._draw_text_logo(slide, x, y, w=Inches(1.0), h=Inches(0.6))

    def _draw_text_logo(self, slide, x, y, w, h):
        # Stylized fallback: 3 small green triangles + wordmark
        for i, off in enumerate([0.0, 0.18, 0.36]):
            tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                         x + Inches(off), y,
                                         Inches(0.22), Inches(0.32))
            tri.fill.solid()
            tri.fill.fore_color.rgb = GREEN_PRIMARY
            tri.line.fill.background()
        _add_text(slide, x, y + Inches(0.34), w, Inches(0.22),
                  "TaomoAI", size=Pt(8), bold=True, color=GREEN_DARK,
                  align=PP_ALIGN.CENTER)

    def draw_footer(self, slide):
        _add_text(slide, Inches(0), FOOTER_Y, SLIDE_W, Inches(0.3),
                  self.footer_text,
                  size=SZ_FOOTER, color=TEXT_MUTED,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # default cover background (industrial cityscape, extracted from template)
    _COVER_BG = os.path.join(ASSETS_DIR, "cover_bg.png")

    # ---- Cover ----
    def draw_title_slide(self, slide, spec, spec_dir):
        """Left half: white + logo + text.
        Right half: spec image → cover_bg.png → dark green panel (fallback).
        Matches the Siemens PPAI Weekly Tracking template exactly.
        """
        half_w = SLIDE_W // 2
        img = _resolve_path(spec_dir, spec.get("image"))
        import os as _os
        # resolve right-side image: explicit spec > cover_bg default > green panel
        right_img = None
        if img and _os.path.isfile(img):
            right_img = img
        elif _os.path.isfile(self._COVER_BG):
            right_img = self._COVER_BG

        if right_img:
            slide.shapes.add_picture(right_img, half_w, Emu(0), half_w, SLIDE_H)
        else:
            _add_rect(slide, half_w, Emu(0), half_w, SLIDE_H, fill=GREEN_DARK)
            if has_logo(LOGO_INDUSTRIALMIND):
                slide.shapes.add_picture(
                    LOGO_INDUSTRIALMIND,
                    half_w + (half_w - Inches(2.5)) // 2, Inches(2.7),
                    width=Inches(2.5))

        # left logo — top-left, matching template
        if has_logo(LOGO_INDUSTRIALMIND):
            slide.shapes.add_picture(LOGO_INDUSTRIALMIND,
                                     MARGIN_X, Inches(0.35),
                                     width=Inches(1.1))
        else:
            self._draw_text_logo(slide, MARGIN_X, Inches(0.35),
                                 Inches(1.2), Inches(0.7))

        # title + subtitle centred vertically on left half — template style
        _add_text(slide, MARGIN_X, Inches(2.7), half_w - MARGIN_X - Inches(0.3),
                  Inches(1.1), spec.get("title", ""),
                  size=SZ_COVER_TITLE, bold=True, color=TEXT_DARK)
        if spec.get("subtitle"):
            _add_text(slide, MARGIN_X, Inches(3.9), half_w - MARGIN_X - Inches(0.3),
                      Inches(0.7), spec["subtitle"],
                      size=Pt(26), bold=True, color=TEXT_DARK)
        if spec.get("date"):
            _add_text(slide, MARGIN_X, Inches(4.75), half_w - MARGIN_X - Inches(0.3),
                      Inches(0.4), spec["date"],
                      size=SZ_SUBTITLE, color=TEXT_MUTED)
        # footer
        _add_text(slide, MARGIN_X, FOOTER_Y, half_w - MARGIN_X, Inches(0.3),
                  "Copyright @ TaomoAI",
                  size=SZ_FOOTER, color=TEXT_MUTED, align=PP_ALIGN.LEFT)

    def draw_section_slide(self, slide, spec):
        _add_text(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(1.2),
                  spec.get("title", "SECTION"),
                  size=SZ_SECTION, bold=True, color=TEXT_DARK,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if spec.get("subtitle"):
            _add_text(slide, Inches(0), Inches(4.4), SLIDE_W, Inches(0.6),
                      spec["subtitle"],
                      size=SZ_SUBTITLE, color=TEXT_MUTED,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    def draw_thanks_slide(self, slide, spec, spec_dir):
        """Mirrors the cover layout exactly — same right-side image, same text chrome."""
        half_w = SLIDE_W // 2
        img = _resolve_path(spec_dir, spec.get("image"))
        import os as _os
        right_img = None
        if img and _os.path.isfile(img):
            right_img = img
        elif _os.path.isfile(self._COVER_BG):
            right_img = self._COVER_BG

        if right_img:
            slide.shapes.add_picture(right_img, half_w, Emu(0), half_w, SLIDE_H)
        else:
            _add_rect(slide, half_w, Emu(0), half_w, SLIDE_H, fill=GREEN_DARK)
            if has_logo(LOGO_INDUSTRIALMIND):
                slide.shapes.add_picture(
                    LOGO_INDUSTRIALMIND,
                    half_w + (half_w - Inches(2.5)) // 2, Inches(2.7),
                    width=Inches(2.5))

        # left logo
        if has_logo(LOGO_INDUSTRIALMIND):
            slide.shapes.add_picture(LOGO_INDUSTRIALMIND,
                                     MARGIN_X, Inches(0.35),
                                     width=Inches(1.1))
        else:
            self._draw_text_logo(slide, MARGIN_X, Inches(0.35),
                                 Inches(1.2), Inches(0.7))

        # main title (bold, same weight as cover title)
        _add_text(slide, MARGIN_X, Inches(2.7), half_w - MARGIN_X - Inches(0.3),
                  Inches(1.1), spec.get("title", "感谢聆听"),
                  size=SZ_COVER_TITLE, bold=True, color=TEXT_DARK)
        if spec.get("tagline"):
            _add_text(slide, MARGIN_X, Inches(3.9), half_w - MARGIN_X - Inches(0.3),
                      Inches(0.6), spec["tagline"],
                      size=Pt(22), bold=True, color=TEXT_DARK)
        if spec.get("contact"):
            _add_text(slide, MARGIN_X, Inches(4.6), half_w - MARGIN_X - Inches(0.3),
                      Inches(0.4), spec["contact"],
                      size=SZ_SUBTITLE, color=TEXT_MUTED)
        # footer
        _add_text(slide, MARGIN_X, FOOTER_Y, half_w - MARGIN_X, Inches(0.3),
                  "Copyright @ TaomoAI",
                  size=SZ_FOOTER, color=TEXT_MUTED, align=PP_ALIGN.LEFT)
