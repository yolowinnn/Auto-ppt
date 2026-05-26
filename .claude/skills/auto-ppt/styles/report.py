"""
Style: report — Comprehensive assessment / evaluation report.

Visual signature (from Data/100-图纸识别模型综合评估报告-0330.pdf):
  - Cover: white background; thin GREEN vertical bar (left) + thin BLUE
    vertical bar (right); large bold title; gray pill subtitle; bright
    green INDUSTRIALMIND.AI wordmark; full-width green bottom bar.
  - Content: green section title top-left; bold-black headline with
    green vertical bar accent; light-gray callout box for key insight;
    KPI cards / data tables; full-width green bottom bar.
  - Right-top: small IndustrialMind logo with "Taomo.ai" wordmark.
"""
import os, sys
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from theme import (
    SLIDE_W, SLIDE_H, MARGIN_X, TITLE_BAR_Y, FOOTER_Y,
    GREEN_PRIMARY, GREEN_DARK, GREEN_BRIGHT, GRAY_DIVIDER, GRAY_CALLOUT,
    TEXT_DARK, TEXT_BODY, TEXT_MUTED, WHITE, KPI_INFO,
    SZ_TITLE, SZ_SUBTITLE, SZ_BODY, SZ_FOOTER, SZ_COVER_TITLE, SZ_SECTION,
    LOGO_INDUSTRIALMIND, LOGO_TAOMO, has_logo,
)
from builder import _add_text, _add_rect, _add_line, _add_image_or_placeholder, _resolve_path


class ReportStyle:
    name = "report"
    bg_color = WHITE
    table_header_bg = RGBColor(0xF1, 0xF8, 0xE9)
    table_header_fg = GREEN_DARK
    accent_color = GREEN_BRIGHT
    bullet_color = GREEN_BRIGHT
    heading_color = GREEN_DARK
    default_footer = ""  # report uses a green bottom bar instead of text

    def __init__(self, meta):
        self.meta = meta or {}
        self.footer_text = self.meta.get("footer", self.default_footer)

    def draw_header(self, slide, *, title, subtitle=None):
        # green section title (smaller than other styles, bold)
        _add_text(slide, MARGIN_X, Inches(0.35),
                  Inches(11), Inches(0.5), title,
                  size=Pt(22), bold=True, color=GREEN_BRIGHT)
        if subtitle:
            # bold-black headline with thin green vertical bar accent
            bar_x = MARGIN_X
            bar_y = Inches(0.95)
            _add_rect(slide, bar_x, bar_y, Pt(3), Inches(0.35),
                      fill=GREEN_BRIGHT)
            _add_text(slide, bar_x + Inches(0.15), bar_y - Inches(0.02),
                      Inches(11), Inches(0.4), subtitle,
                      size=Pt(15), bold=True, color=TEXT_DARK)

    def draw_logo(self, slide):
        x = SLIDE_W - Inches(1.4)
        y = Inches(0.25)
        if has_logo(LOGO_TAOMO):
            slide.shapes.add_picture(LOGO_TAOMO, x, y, width=Inches(1.0))
        elif has_logo(LOGO_INDUSTRIALMIND):
            slide.shapes.add_picture(LOGO_INDUSTRIALMIND, x, y,
                                     width=Inches(1.0))
        else:
            self._draw_text_logo(slide, x, y)

    def _draw_text_logo(self, slide, x, y, color=GREEN_PRIMARY):
        for i, off in enumerate([0.0, 0.18, 0.36]):
            tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                         x + Inches(off), y,
                                         Inches(0.22), Inches(0.32))
            tri.fill.solid()
            tri.fill.fore_color.rgb = color
            tri.line.fill.background()
        _add_text(slide, x - Inches(0.05), y + Inches(0.34), Inches(0.7),
                  Inches(0.22), "Taomo.ai",
                  size=Pt(8), bold=True, color=color, align=PP_ALIGN.CENTER)

    def draw_footer(self, slide):
        # full-width thin green bottom bar
        _add_rect(slide, Inches(0), SLIDE_H - Inches(0.2),
                  SLIDE_W, Inches(0.2), fill=GREEN_PRIMARY)
        if self.footer_text:
            _add_text(slide, MARGIN_X, FOOTER_Y - Inches(0.05),
                      SLIDE_W - 2 * MARGIN_X, Inches(0.3),
                      self.footer_text, size=SZ_FOOTER, color=WHITE,
                      align=PP_ALIGN.CENTER)

    # ---- Cover (side bars + bottom bar) ----
    def draw_title_slide(self, slide, spec, spec_dir):
        # thin green vertical bar on the left
        _add_rect(slide, Inches(2.0), Inches(1.3), Pt(3.5),
                  Inches(4.5), fill=GREEN_BRIGHT)
        # thin blue vertical bar on the right
        _add_rect(slide, SLIDE_W - Inches(2.0), Inches(1.3),
                  Pt(3.5), Inches(4.5), fill=KPI_INFO)
        # logo top-right
        if has_logo(LOGO_INDUSTRIALMIND):
            slide.shapes.add_picture(LOGO_INDUSTRIALMIND,
                                     SLIDE_W - Inches(1.4), Inches(0.3),
                                     width=Inches(1.0))
        else:
            self._draw_text_logo(slide, SLIDE_W - Inches(1.4), Inches(0.3))
        # title (large, bold black)
        _add_text(slide, Inches(2.3), Inches(2.0),
                  SLIDE_W - Inches(4.6), Inches(1.0),
                  spec.get("title", ""),
                  size=Pt(34), bold=True, color=TEXT_DARK,
                  align=PP_ALIGN.CENTER)
        if spec.get("subtitle"):
            _add_text(slide, Inches(2.3), Inches(3.0),
                      SLIDE_W - Inches(4.6), Inches(0.6),
                      spec["subtitle"],
                      size=Pt(20), color=TEXT_DARK, align=PP_ALIGN.CENTER)
        # gray pill containing scope/dataset/date
        if spec.get("scope"):
            pill_w = Inches(7)
            pill_x = (SLIDE_W - pill_w) // 2
            _add_rect(slide, pill_x, Inches(3.9), pill_w, Inches(0.45),
                      fill=GRAY_CALLOUT)
            _add_text(slide, pill_x, Inches(3.9), pill_w, Inches(0.45),
                      spec["scope"], size=Pt(13), color=TEXT_BODY,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # bright green wordmark
        _add_text(slide, Inches(2.3), Inches(4.6),
                  SLIDE_W - Inches(4.6), Inches(0.5),
                  spec.get("wordmark", "INDUSTRIALMIND.AI"),
                  size=Pt(18), bold=True, color=GREEN_BRIGHT,
                  align=PP_ALIGN.CENTER)
        # bottom green bar
        _add_rect(slide, Inches(0), SLIDE_H - Inches(0.6),
                  SLIDE_W, Inches(0.6), fill=GREEN_PRIMARY)

    def draw_section_slide(self, slide, spec):
        _add_text(slide, Inches(0), Inches(3.0), SLIDE_W, Inches(1.2),
                  spec.get("title", "SECTION"),
                  size=SZ_SECTION, bold=True, color=GREEN_DARK,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if spec.get("subtitle"):
            _add_text(slide, Inches(0), Inches(4.3), SLIDE_W, Inches(0.6),
                      spec["subtitle"], size=SZ_SUBTITLE, color=TEXT_MUTED,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_rect(slide, Inches(0), SLIDE_H - Inches(0.2),
                  SLIDE_W, Inches(0.2), fill=GREEN_PRIMARY)

    def draw_thanks_slide(self, slide, spec, spec_dir):
        _add_rect(slide, Inches(2.0), Inches(1.3), Pt(3.5),
                  Inches(4.5), fill=GREEN_BRIGHT)
        _add_rect(slide, SLIDE_W - Inches(2.0), Inches(1.3),
                  Pt(3.5), Inches(4.5), fill=KPI_INFO)
        _add_text(slide, MARGIN_X, Inches(3.0), SLIDE_W - 2 * MARGIN_X,
                  Inches(1.0), spec.get("title", "Thank you"),
                  size=SZ_COVER_TITLE, bold=True, color=TEXT_DARK,
                  align=PP_ALIGN.CENTER)
        if spec.get("subtitle"):
            _add_text(slide, MARGIN_X, Inches(4.0), SLIDE_W - 2 * MARGIN_X,
                      Inches(0.5), spec["subtitle"],
                      size=SZ_SUBTITLE, color=TEXT_MUTED,
                      align=PP_ALIGN.CENTER)
        _add_text(slide, MARGIN_X, Inches(4.7), SLIDE_W - 2 * MARGIN_X,
                  Inches(0.5),
                  spec.get("wordmark", "INDUSTRIALMIND.AI"),
                  size=Pt(18), bold=True, color=GREEN_BRIGHT,
                  align=PP_ALIGN.CENTER)
        _add_rect(slide, Inches(0), SLIDE_H - Inches(0.6),
                  SLIDE_W, Inches(0.6), fill=GREEN_PRIMARY)
